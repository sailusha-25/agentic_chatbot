# core/file_handler.py

import os
import pandas as pd
from langchain.text_splitter import RecursiveCharacterTextSplitter
import pypdf
import docx
import pptx

def parse_document(uploaded_file):
    """
    Parses an uploaded file based on its extension and extracts text.

    Args:
        uploaded_file: A file-like object from Streamlit's file_uploader.

    Returns:
        str: The extracted text content from the document.
    """
    text = ""
    file_name = uploaded_file.name
    _, extension = os.path.splitext(file_name)

    if extension == ".pdf":
        pdf_reader = pypdf.PdfReader(uploaded_file)
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
    elif extension == ".docx":
        doc = docx.Document(uploaded_file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif extension == ".pptx":
        prs = pptx.Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif extension == ".csv":
        df = pd.read_csv(uploaded_file)
        text = df.to_string()
    elif extension in [".txt", ".md"]:
        text = uploaded_file.getvalue().decode("utf-8")

    return text

def get_text_chunks(text):
    """
    Splits a long text into smaller chunks for processing.

    Args:
        text (str): The input text.

    Returns:
        list[str]: A list of text chunks.
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200
    )
    chunks = text_splitter.split_text(text)
    return chunks